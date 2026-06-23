"""PPTX document loader — extract slides as text with metadata."""

import logging
from pathlib import Path

from pptx import Presentation

from src.exceptions import DocumentLoadError
from src.models import SourceChunk

logger = logging.getLogger(__name__)


def load_pptx(path: Path) -> list[SourceChunk]:
    """Load a PPTX file and return slides as SourceChunks.

    Args:
        path: Path to the PPTX file.

    Returns:
        List of SourceChunk, one per slide.

    Raises:
        DocumentLoadError: If the PPTX cannot be read.
    """
    try:
        prs = Presentation(str(path))
        chunks: list[SourceChunk] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)

            if text_parts:
                chunks.append(
                    SourceChunk(
                        content="\n".join(text_parts),
                        source=path.name,
                        page=slide_num,
                    )
                )

        logger.info("Loaded %d slides from %s", len(chunks), path.name)
        return chunks

    except Exception as e:
        raise DocumentLoadError(str(path), str(e)) from e


def load_full_pptx(path: Path) -> str:
    """Load entire PPTX as a single text string (for artifacts).

    Args:
        path: Path to the PPTX file.

    Returns:
        Full text content of all slides.

    Raises:
        DocumentLoadError: If the PPTX cannot be read.
    """
    try:
        prs = Presentation(str(path))
        slides = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)

            if text_parts:
                slides.append(f"[Slide {slide_num}]\n" + "\n".join(text_parts))

        return "\n\n".join(slides)

    except Exception as e:
        raise DocumentLoadError(str(path), str(e)) from e
